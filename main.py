MODULE_NAME = 'module_name'
UNIT_NAME = 'unit_name'
GOALS = 'goals'
FOR_WHAT = 'for_what'
ROUTE = 'route'

unit = {
    MODULE_NAME: 'Основы программирования',
    UNIT_NAME: 'Git и Github. Базовые команды для работы',
    GOALS: [
        'Создавать репозитории на github',
        'Пользоваться git для контроля версий кода',
        'Пользоваться github для хранения кода'
    ],
    FOR_WHAT: [
        'Использовать git и github для профессиональной разработки программ'
    ],
    ROUTE: [
        'Системы контроля версий',
        'Git',
        'Github',
        'Основные команды git'
    ]
}

import copy
import six


def duplicate_slide(pres, index):
    template = pres.slides[index]
    # try:
    #     blank_slide_layout = pres.slide_layouts[12]
    # except:
    #     blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)-1]
    blank_slide_layout = pres.slide_layouts[0]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )

    return copied_slide


# Надо из unit сделать презентацию :)
from pptx import Presentation

prs = Presentation('template.pptx')

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for number, slide in enumerate(prs.slides, 0):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for k, v in unit.items():
                    if k in run.text:
                        if 'item' in run.text:
                            for text_item in unit[k][1:]:
                                # Меняем текст слайда
                                run.text = run.text.replace(f'{k}.item', text_item)
                                # Вставляет в конец и то хорошо :)
                                duplicate_slide(pres=prs, index=number)
                                # обратно меняем на item
                                run.text = f'{k}.item'
                            run.text = run.text.replace(f'{k}.item', unit[k][0])
                            # if 'item.' in run.text:
                            #     pass
                            # else:
                            #     # Значит надо размножить этот слайд
                            #     for i, text_item in enumerate(unit[k], 0):
                            #         # Меняем текст слайда
                            #         run.text = f'run.text.{i}'
                            #         # Вставляет в конец и то хорошо :)
                            #         duplicate_slide(pres=prs, index=number)
                        else:
                            if isinstance(v, str):
                                run.text = run.text.replace(k, v)
                            elif isinstance(v, list):
                                # TODO: работает коряво, список едет
                                text = '\n'.join(v)
                                run.text = run.text.replace(k, text)
                # text_runs.append(run.text)

for text in text_runs:
    print(text)

prs.save('result.pptx')
